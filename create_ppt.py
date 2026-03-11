from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_pitch_deck():
    prs = Presentation()

    # Define some colors and styles
    NAVY = RGBColor(0, 32, 96)
    GOLD = RGBColor(191, 144, 0)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background image
    hero_img = r"C:\Users\ravim\.gemini\antigravity\brain\f1011572-d5bc-4d4e-b7f4-53f8d8cc21f8\pitch_deck_hero_1773168994497.png"
    if os.path.exists(hero_img):
        prs.slides[0].shapes.add_picture(hero_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Title Text
    left = top = Inches(2)
    width = prs.slide_width - Inches(4)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Art&Auction"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Discover, Bid, and Win Amazing Items Across India"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Trust Gap in Indian Auctions"
    content = slide.placeholders[1].text_frame
    content.text = "• Fragmented Market: Highly fragmented niche auctions.\n• Trust Deficit: Lack of verified sellers and records.\n• Currency Friction: No native ₹ INR optimization.\n• Data Loss: Item records often disappear on smaller platforms."

    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution: A Modern, Robust Platform"
    content = slide.placeholders[1].text_frame
    content.text = "• Localized Experience: Built for the Indian market (₹).\n• Immutable Persistence: Vercel Blob ensures zero data loss.\n• Smart Bidding: Automated status tracking and increments.\n• Verified Community: Secure role-based access control."

    # Slide 4: Market Opportunity
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title and Content
    slide.shapes.title.text = "Market Opportunity: Indian Collectibles"
    left = Inches(0.5); top = Inches(1.5); width = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph(); p.text = "• 25% YoY growth in online auctions in India.\n• Growing demographic of young collectors.\n• High-value transactions moving online."
    
    chart_img = r"C:\Users\ravim\.gemini\antigravity\brain\f1011572-d5bc-4d4e-b7f4-53f8d8cc21f8\market_growth_chart_1773169030628.png"
    if os.path.exists(chart_img):
        slide.shapes.add_picture(chart_img, Inches(5.5), Inches(2), width=Inches(3.5))

    # Slide 5: Product Showcase
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Seamless Bidding Everywhere"
    mockup_img = r"C:\Users\ravim\.gemini\antigravity\brain\f1011572-d5bc-4d4e-b7f4-53f8d8cc21f8\product_showcase_mockup_1773169011749.png"
    if os.path.exists(mockup_img):
        slide.shapes.add_picture(mockup_img, Inches(0.5), Inches(1.5), height=Inches(5))
    
    left = Inches(6); top = Inches(2); width = Inches(3.5)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph(); p.text = "• Real-time Dashboards\n• Live Notification Engine\n• 60-second Listing Flow"

    # Slide 6: Technical Excellence
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technical Excellence"
    content = slide.placeholders[1].text_frame
    content.text = "• Tech Stack: Python/Flask + SQLite + Vercel Blob.\n• Reliability: Synchronous cloud synchronization.\n• UI/UX: Premium Dark mode + AOS animations.\n• SEO: Structured metadata for Indian buyers."

    # Slide 7: Business Advantage
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Advantage"
    content = slide.placeholders[1].text_frame
    content.text = "• Zero Maintenance: 100% Cloud native hosting.\n• Trust Shield: Verified deletion and records management.\n• Efficiency: Low cost of execution per transaction."

    # Slide 8: Contact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Join Us in Digitizing Indian Auctions"
    content = slide.placeholders[1].text_frame
    content.text = "Website: artnauction-one.vercel.app\nEmail: hi@artnauction.in\nSocials: @ArtAndAuction_IN"

    prs.save('Art_and_Auction_Pitch_Deck.pptx')
    print("Pitch Deck saved as Art_and_Auction_Pitch_Deck.pptx")

if __name__ == "__main__":
    create_pitch_deck()
